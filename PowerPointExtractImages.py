#!/usr/bin/env python3

import sys
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images_from_shape(shape, slide_number, shape_label, image_counter):
	if hasattr(shape, "image"):
		image_filename = f"{str(image_counter).zfill(3)}-slide-{slide_number}-shape-{shape_label}.jpeg"
		with open(image_filename, "wb") as f:
			f.write(shape.image.blob)
		print(f"\t\tExtracted image to {image_filename}")
		image_counter += 1

	# If the shape is a GROUP, recurse into its .shapes
	if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
		for idx, subshape in enumerate(shape.shapes, start=1):
			# Build a label such as "3-1", "3-2" for subshapes
			subshape_label = f"{shape_label}-{idx}"
			image_counter = extract_images_from_shape(subshape, slide_number, subshape_label, image_counter)

	return image_counter


def main():
	if len(sys.argv) < 2:
		print("Usage: python extract_images.py <YourPresentation.pptx>")
		sys.exit(1)

	pptx_file = sys.argv[1]
	presentation = pptx.Presentation(pptx_file)

	image_counter = 1
	for slide_idx, slide in enumerate(presentation.slides, start=1):
		print(f"Processing Slide [{slide_idx}] of {len(presentation.slides)}")
		for shape_idx, shape in enumerate(slide.shapes, start=1):
			print(f"\tProcessing Shape [{shape_idx}] of {len(slide.shapes)}")
			image_counter = extract_images_from_shape(
				shape,
				slide_idx,
				str(shape_idx),
				image_counter
			)

if __name__ == "__main__":
	main()